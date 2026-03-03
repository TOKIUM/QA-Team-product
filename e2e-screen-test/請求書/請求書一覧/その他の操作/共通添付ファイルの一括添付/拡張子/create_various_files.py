"""
各種拡張子のファイルおよび名前がないファイル（拡張子なし）を作成するスクリプト

作成対象:
  .jpg / .jpeg / .png / .pdf / .xlsx / .xls / .csv / .txt
  .doc / .docx / .gif / .pptx / 名前なしファイル
"""

import os
import struct
from io import BytesIO

from PIL import Image
from openpyxl import Workbook
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))


# ─────────────────────────────────────────────
# 個別ファイル作成関数
# ─────────────────────────────────────────────
def create_jpg(path: str):
    """1x1 白画像の JPEG"""
    img = Image.new("RGB", (1, 1), "white")
    img.save(path, "JPEG")


def create_jpeg(path: str):
    """1x1 白画像の JPEG (.jpeg 拡張子)"""
    img = Image.new("RGB", (1, 1), "white")
    img.save(path, "JPEG")


def create_png(path: str):
    """1x1 白画像の PNG"""
    img = Image.new("RGBA", (1, 1), "white")
    img.save(path, "PNG")


def create_pdf(path: str):
    """1ページの PDF"""
    c = canvas.Canvas(path, pagesize=A4)
    c.setFont("Helvetica", 24)
    c.drawString(72, 750, "Sample PDF")
    c.save()


def create_xlsx(path: str):
    """セルにデータを持つ xlsx"""
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Sample XLSX"
    wb.save(path)


def create_xls(path: str):
    """
    最小限の BIFF8 (Excel 97-2003) バイナリを直接生成する。
    外部ライブラリ不要。空の BOF + EOF のみ。
    """
    bof_record = struct.pack("<HH", 0x0809, 8)  # BOF record type, data length
    bof_data = struct.pack("<HHHH", 0x0600, 0x0010, 0x0000, 0x0000)
    eof_record = struct.pack("<HH", 0x000A, 0)  # EOF record
    with open(path, "wb") as f:
        f.write(bof_record + bof_data + eof_record)


def create_csv(path: str):
    """UTF-8 BOM 付き CSV"""
    with open(path, "w", encoding="utf-8-sig", newline="") as f:
        f.write("Column1,Column2,Column3\n")
        f.write("A,B,C\n")


def create_txt(path: str):
    """UTF-8 テキスト"""
    with open(path, "w", encoding="utf-8") as f:
        f.write("Sample text file.\n")


def create_doc(path: str):
    """
    最小限の Word 97-2003 (.doc) バイナリ。
    CFB (Compound File Binary) の簡易ヘッダのみ。
    """
    header = bytearray(512)
    # CFB signature
    header[0:8] = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
    # Minor version = 0x003E, Major version = 0x0003
    struct.pack_into("<H", header, 0x18, 0x003E)
    struct.pack_into("<H", header, 0x1A, 0x0003)
    # Byte order little-endian
    struct.pack_into("<H", header, 0x1C, 0xFFFE)
    # Sector size power = 9 (512 bytes)
    struct.pack_into("<H", header, 0x1E, 0x0009)
    # Mini sector size power = 6 (64 bytes)
    struct.pack_into("<H", header, 0x20, 0x0006)
    # Fill DIFAT with 0xFFFFFFFF
    for i in range(0x4C, 0x200, 4):
        struct.pack_into("<I", header, i, 0xFFFFFFFF)
    # End of chain markers
    struct.pack_into("<I", header, 0x2C, 0xFFFFFFFE)  # First directory sector SECT
    struct.pack_into("<I", header, 0x3C, 0xFFFFFFFE)  # First mini FAT sector
    struct.pack_into("<I", header, 0x44, 0xFFFFFFFE)  # First DIFAT sector
    with open(path, "wb") as f:
        f.write(header)


def create_docx(path: str):
    """python-docx で docx 作成"""
    doc = Document()
    doc.add_paragraph("Sample DOCX")
    doc.save(path)


def create_gif(path: str):
    """1x1 白画像の GIF"""
    img = Image.new("P", (1, 1), 0)
    img.save(path, "GIF")


def create_pptx(path: str):
    """python-pptx で pptx 作成"""
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.paragraphs[0].text = "Sample PPTX"
    txBox.text_frame.paragraphs[0].font.size = Pt(24)
    prs.save(path)


def create_no_name(path: str):
    """名前がないファイル（拡張子なし、ファイル名のみ空に近い形）"""
    with open(path, "w", encoding="utf-8") as f:
        f.write("This file has no name (no extension, minimal name).\n")


# ─────────────────────────────────────────────
# ファイル定義: (ファイル名, 説明, 作成関数)
# ─────────────────────────────────────────────
FILE_DEFINITIONS = [
    ("sample.jpg", ".jpg", create_jpg),
    ("sample.jpeg", ".jpeg", create_jpeg),
    ("sample.png", ".png", create_png),
    ("sample.pdf", ".pdf", create_pdf),
    ("sample.xlsx", ".xlsx", create_xlsx),
    ("sample.xls", ".xls", create_xls),
    ("sample.csv", ".csv", create_csv),
    ("sample.txt", ".txt", create_txt),
    ("sample.doc", ".doc", create_doc),
    ("sample.docx", ".docx", create_docx),
    ("sample.gif", ".gif", create_gif),
    ("sample.pptx", ".pptx", create_pptx),
    ("_", "名前なしファイル", create_no_name),
]


def main():
    success = 0
    fail = 0

    for filename, desc, creator_func in FILE_DEFINITIONS:
        filepath = os.path.join(OUTPUT_DIR, filename)
        print(f"Creating: {filename:20s} ({desc})")
        try:
            creator_func(filepath)
            if os.path.exists(filepath):
                size = os.path.getsize(filepath)
                print(f"  -> OK ({size:,} bytes)")
                success += 1
            else:
                print("  -> FAIL: file not found after creation")
                fail += 1
        except Exception as e:
            print(f"  -> ERROR: {e}")
            fail += 1

    print(f"\nDone: {success} succeeded, {fail} failed (total {success + fail})")


if __name__ == "__main__":
    main()
